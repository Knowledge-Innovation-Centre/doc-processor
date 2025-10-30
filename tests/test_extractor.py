"""
Tests for ContentExtractor.
"""

from pathlib import Path

import pytest

from docprocessor.core.extractor import ContentExtractionError, ContentExtractor


class TestContentExtractor:
    """Tests for ContentExtractor class."""

    def test_init(self):
        """Test extractor initialization."""
        extractor = ContentExtractor()

        assert ".pdf" in extractor.supported_extensions
        assert ".txt" in extractor.supported_extensions
        assert ".docx" in extractor.supported_extensions
        assert ".md" in extractor.supported_extensions

    def test_extract_txt_file(self, sample_txt_file):
        """Test extracting text from .txt file."""
        extractor = ContentExtractor()

        result = extractor.extract(sample_txt_file)

        assert "text" in result
        assert "page_count" in result
        assert "metadata" in result
        assert len(result["text"]) > 0
        assert result["page_count"] == 1
        assert result["metadata"]["format"] == "txt"

    def test_extract_md_file(self, tmp_path):
        """Test extracting text from .md file."""
        md_file = tmp_path / "test.md"
        md_file.write_text("# Heading\n\nThis is markdown content.")

        extractor = ContentExtractor()
        result = extractor.extract(md_file)

        assert "text" in result
        assert "markdown" in result["text"].lower() or "heading" in result["text"].lower()
        assert result["metadata"]["format"] == "md"

    def test_extract_nonexistent_file(self):
        """Test extracting from non-existent file raises error."""
        extractor = ContentExtractor()

        with pytest.raises(ContentExtractionError, match="File not found"):
            extractor.extract(Path("/nonexistent/file.txt"))

    def test_extract_with_unicode(self, tmp_path):
        """Test extracting file with unicode characters."""
        unicode_file = tmp_path / "unicode.txt"
        unicode_file.write_text("Hello 世界 🌍", encoding="utf-8")

        extractor = ContentExtractor()
        result = extractor.extract(unicode_file)

        assert "Hello" in result["text"]
        assert "世界" in result["text"]

    def test_extract_latin1_encoding(self, tmp_path):
        """Test extracting file with latin-1 encoding."""
        latin1_file = tmp_path / "latin1.txt"
        latin1_file.write_bytes("Café résumé".encode("latin-1"))

        extractor = ContentExtractor()
        result = extractor.extract(latin1_file)

        assert len(result["text"]) > 0

    def test_is_supported(self):
        """Test is_supported method."""
        extractor = ContentExtractor()

        assert extractor.is_supported(".pdf") is True
        assert extractor.is_supported(".txt") is True
        assert extractor.is_supported(".docx") is True
        assert extractor.is_supported(".pptx") is True
        assert extractor.is_supported(".xyz") is False

    def test_extract_unsupported_file(self, tmp_path):
        """Test extracting unsupported file type."""
        unsupported_file = tmp_path / "test.xyz"
        unsupported_file.write_text("test content")

        extractor = ContentExtractor()
        result = extractor.extract(unsupported_file)

        # Should fall back to text extraction
        assert "text" in result
        assert result["text"] == "test content"

    def test_extract_empty_file(self, tmp_path):
        """Test extracting empty file."""
        empty_file = tmp_path / "empty.txt"
        empty_file.write_text("")

        extractor = ContentExtractor()
        result = extractor.extract(empty_file)

        assert result["text"] == ""
        assert result["page_count"] == 1

    def test_extract_large_file(self, tmp_path):
        """Test extracting large text file."""
        large_file = tmp_path / "large.txt"
        large_content = "Line of text.\n" * 10000
        large_file.write_text(large_content)

        extractor = ContentExtractor()
        result = extractor.extract(large_file)

        assert len(result["text"]) > 100000
        assert result["page_count"] == 1

    def test_extract_docx_not_installed(self, tmp_path, monkeypatch):
        """Test DOCX extraction when python-docx is not installed."""
        docx_file = tmp_path / "test.docx"
        docx_file.touch()

        extractor = ContentExtractor()

        # Mock docx import to raise ImportError
        import builtins

        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "docx":
                raise ImportError("No module named 'docx'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)

        with pytest.raises(ContentExtractionError, match="python-docx not installed"):
            extractor._extract_docx(docx_file)

    def test_extract_pdf_file(self, pdf_with_text_content):
        """Test extracting text from PDF file."""
        extractor = ContentExtractor()

        result = extractor.extract(pdf_with_text_content)

        assert "text" in result
        assert "page_count" in result
        assert result["metadata"]["format"] == "pdf"
        assert result["metadata"]["extraction_method"] == "ocr_pipeline"
        assert len(result["text"]) > 0

    def test_extract_docx_with_content(self, tmp_path):
        """Test extracting text from DOCX with paragraphs and tables."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        # Create DOCX with content
        docx_file = tmp_path / "test.docx"
        doc = Document()
        doc.add_paragraph("First paragraph")
        doc.add_paragraph("Second paragraph")

        # Add table
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Data 1"
        table.cell(1, 1).text = "Data 2"

        doc.save(str(docx_file))

        extractor = ContentExtractor()
        result = extractor.extract(docx_file)

        assert "text" in result
        assert "First paragraph" in result["text"]
        assert "Second paragraph" in result["text"]
        assert result["metadata"]["format"] == "docx"
        assert result["metadata"]["extraction_method"] == "python-docx"
        assert result["metadata"]["paragraph_count"] == 2
        assert result["metadata"]["table_count"] == 1

    def test_extract_docx_with_empty_paragraphs(self, tmp_path):
        """Test extracting DOCX skips empty paragraphs."""
        try:
            from docx import Document
        except ImportError:
            pytest.skip("python-docx not installed")

        docx_file = tmp_path / "empty_paras.docx"
        doc = Document()
        doc.add_paragraph("Content")
        doc.add_paragraph("")  # Empty
        doc.add_paragraph("   ")  # Whitespace only
        doc.add_paragraph("More content")
        doc.save(str(docx_file))

        extractor = ContentExtractor()
        result = extractor.extract(docx_file)

        # Should only count non-empty paragraphs
        assert result["metadata"]["paragraph_count"] == 2

    def test_extract_image_file(self, tmp_path):
        """Test extracting text from image using OCR."""
        try:
            import pytesseract
            from PIL import Image

            # Check if tesseract is actually available
            pytesseract.get_tesseract_version()
        except (ImportError, pytesseract.TesseractNotFoundError):
            pytest.skip("PIL, pytesseract or tesseract binary not installed")

        # Create simple test image
        img_file = tmp_path / "test.png"
        img = Image.new("RGB", (100, 30), color="white")
        img.save(img_file)

        extractor = ContentExtractor()
        result = extractor.extract(img_file)

        assert "text" in result
        assert result["metadata"]["format"] == "png"
        assert result["metadata"]["extraction_method"] == "tesseract_ocr"
        assert "image_size" in result["metadata"]

    def test_extract_image_missing_dependencies(self, tmp_path, monkeypatch):
        """Test image extraction fails gracefully without dependencies."""
        # Create a dummy image file
        img_file = tmp_path / "test.jpg"
        img_file.write_bytes(b"\xff\xd8\xff\xe0")  # JPEG header

        # Mock ImportError for PIL
        import builtins

        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name in ("PIL", "pytesseract"):
                raise ImportError(f"No module named '{name}'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)

        extractor = ContentExtractor()

        with pytest.raises(ContentExtractionError, match="PIL and pytesseract required"):
            extractor.extract(img_file)

    def test_extract_docx_corruption(self, tmp_path):
        """Test DOCX extraction handles corrupted files."""
        try:
            import docx  # noqa: F401
        except ImportError:
            pytest.skip("python-docx not installed")

        # Create a corrupted DOCX file
        bad_docx = tmp_path / "corrupt.docx"
        bad_docx.write_bytes(b"not a real docx file")

        extractor = ContentExtractor()

        with pytest.raises(ContentExtractionError, match="Failed to extract"):
            extractor.extract(bad_docx)

    def test_extract_pptx_not_installed(self, tmp_path, monkeypatch):
        """Test PPTX extraction when python-pptx is not installed."""
        pptx_file = tmp_path / "test.pptx"
        pptx_file.touch()

        extractor = ContentExtractor()

        # Mock pptx import to raise ImportError
        import builtins

        original_import = builtins.__import__

        def mock_import(name, *args, **kwargs):
            if name == "pptx":
                raise ImportError("No module named 'pptx'")
            return original_import(name, *args, **kwargs)

        monkeypatch.setattr(builtins, "__import__", mock_import)

        with pytest.raises(ContentExtractionError, match="python-pptx not installed"):
            extractor._extract_pptx(pptx_file)

    def test_extract_pptx_with_content(self, tmp_path):
        """Test extracting text from PPTX with slides and content."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Create PPTX with content
        pptx_file = tmp_path / "test.pptx"
        prs = Presentation()

        # Add slide 1 with title and content
        slide_layout = prs.slide_layouts[1]  # Title and content layout
        slide1 = prs.slides.add_slide(slide_layout)
        slide1.shapes.title.text = "First Slide Title"
        slide1.placeholders[1].text = "This is the content of the first slide."

        # Add slide 2 with different content
        slide2 = prs.slides.add_slide(slide_layout)
        slide2.shapes.title.text = "Second Slide Title"
        slide2.placeholders[1].text = "This is the content of the second slide."

        prs.save(str(pptx_file))

        extractor = ContentExtractor()
        result = extractor.extract(pptx_file)

        assert "text" in result
        assert "First Slide Title" in result["text"]
        assert "Second Slide Title" in result["text"]
        assert "first slide" in result["text"]
        assert "second slide" in result["text"]
        assert result["page_count"] == 2
        assert result["metadata"]["format"] == "pptx"
        assert result["metadata"]["extraction_method"] == "python-pptx"
        assert result["metadata"]["slide_count"] == 2
        assert result["metadata"]["shape_count"] > 0

    def test_extract_pptx_with_tables(self, tmp_path):
        """Test extracting text from PPTX with tables."""
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Create PPTX with table
        pptx_file = tmp_path / "table_test.pptx"
        prs = Presentation()

        # Add slide with blank layout
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add a table
        rows, cols = 3, 2
        left = Inches(2)
        top = Inches(2)
        width = Inches(4)
        height = Inches(2)

        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table

        # Fill table with data
        table.cell(0, 0).text = "Header 1"
        table.cell(0, 1).text = "Header 2"
        table.cell(1, 0).text = "Row 1 Col 1"
        table.cell(1, 1).text = "Row 1 Col 2"
        table.cell(2, 0).text = "Row 2 Col 1"
        table.cell(2, 1).text = "Row 2 Col 2"

        prs.save(str(pptx_file))

        extractor = ContentExtractor()
        result = extractor.extract(pptx_file)

        assert "text" in result
        assert "Header 1" in result["text"]
        assert "Header 2" in result["text"]
        assert "Row 1 Col 1" in result["text"]
        assert result["metadata"]["has_tables"] is True

    def test_extract_pptx_with_notes(self, tmp_path):
        """Test extracting text from PPTX with speaker notes."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Create PPTX with notes
        pptx_file = tmp_path / "notes_test.pptx"
        prs = Presentation()

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide with Notes"

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_text_frame = notes_slide.notes_text_frame
        notes_text_frame.text = "These are important speaker notes for the presentation."

        prs.save(str(pptx_file))

        extractor = ContentExtractor()
        result = extractor.extract(pptx_file)

        assert "text" in result
        assert "Slide with Notes" in result["text"]
        assert "speaker notes" in result["text"]
        assert "[Notes for Slide 1]" in result["text"]

    def test_extract_pptx_empty(self, tmp_path):
        """Test extracting empty PPTX file."""
        try:
            from pptx import Presentation
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Create empty PPTX
        pptx_file = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(str(pptx_file))

        extractor = ContentExtractor()
        result = extractor.extract(pptx_file)

        assert "text" in result
        assert result["page_count"] == 0
        assert result["metadata"]["slide_count"] == 0

    def test_extract_pptx_corruption(self, tmp_path):
        """Test PPTX extraction handles corrupted files."""
        try:
            import pptx  # noqa: F401
        except ImportError:
            pytest.skip("python-pptx not installed")

        # Create a corrupted PPTX file
        bad_pptx = tmp_path / "corrupt.pptx"
        bad_pptx.write_bytes(b"not a real pptx file")

        extractor = ContentExtractor()

        with pytest.raises(ContentExtractionError, match="Failed to extract"):
            extractor.extract(bad_pptx)


class TestContentExtractionError:
    """Tests for ContentExtractionError exception."""

    def test_exception_creation(self):
        """Test creating ContentExtractionError."""
        error = ContentExtractionError("Test error message")

        assert str(error) == "Test error message"
        assert isinstance(error, Exception)

    def test_exception_raising(self):
        """Test raising ContentExtractionError."""
        with pytest.raises(ContentExtractionError, match="Test error"):
            raise ContentExtractionError("Test error")
